"""
Presentación corporativa AndroidCommonDoc — rediseño con layout de dos columnas.
Columna izquierda: bullets. Columna derecha: visuales reales (tarjetas, diagramas, gráficos).
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta ────────────────────────────────────────────────────────────────────
BG         = RGBColor(0x0D, 0x1B, 0x2A)
CARD       = RGBColor(0x16, 0x2B, 0x3E)
CARD2      = RGBColor(0x1E, 0x38, 0x52)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xCC, 0xD6, 0xE0)
GREY2      = RGBColor(0x88, 0xA0, 0xB0)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
ORANGE     = RGBColor(0xF3, 0x96, 0x12)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW     = RGBColor(0xF1, 0xC4, 0x0F)
PURPLE     = RGBColor(0x9B, 0x59, 0xB6)

W = 13.33
H = 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# ── Primitivas ────────────────────────────────────────────────────────────────
def R(slide, l, t, w, h, fill, line_color=None, line_w=0.5, radius=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def T(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def arrow_right(slide, l, t, w, h, color):
    """Flecha apuntando a la derecha."""
    shape = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def kpi_card(slide, l, t, w, h, number, label, color=ACCENT):
    R(slide, l, t, w, h, CARD2, line_color=color, line_w=1.5)
    R(slide, l, t, w, 0.06, color)
    T(slide, number, l, t+0.08, w, 0.6, size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(slide, label,  l, t+0.6, w, 0.5, size=13, color=GREY, align=PP_ALIGN.CENTER)

def tier_box(slide, l, t, w, h, label, sublabel, color, fill=None):
    R(slide, l, t, w, h, fill or CARD2, line_color=color, line_w=2)
    R(slide, l, t, 0.08, h, color)
    T(slide, label,    l+0.15, t+0.10, w-0.2, 0.38, size=20, bold=True, color=color)
    T(slide, sublabel, l+0.15, t+0.48, w-0.2, 0.45, size=12, color=GREY)

def step_pill(slide, l, t, w, h, num, label, color=ACCENT):
    R(slide, l, t, w, h, CARD2, line_color=color, line_w=1.2)
    R(slide, l, t, 0.36, h, color)
    T(slide, str(num), l+0.05, t+0.06, 0.26, h-0.1, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(slide, label,    l+0.42, t+0.08, w-0.48, h-0.15, size=14, color=GREY)

def flow_box(slide, l, t, w, h, label, color=ACCENT):
    R(slide, l, t, w, h, CARD2, line_color=color, line_w=1.2)
    T(slide, label, l, t, w, h, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

def bar(slide, l, t, w, h, label, value_text, color, pct):
    """Barra horizontal con etiqueta."""
    T(slide, label, l, t, 2.1, 0.28, size=12, color=GREY)
    R(slide, l, t+0.3, w, 0.22, CARD2)
    R(slide, l, t+0.3, w * pct, 0.22, color)
    T(slide, value_text, l + w*pct + 0.05, t+0.28, 1.0, 0.26, size=12, bold=True, color=color)

def rule_chip(slide, l, t, w, h, name, desc, color):
    R(slide, l, t, w, h, CARD2, line_color=color, line_w=1)
    R(slide, l, t, w, 0.05, color)
    T(slide, name, l+0.08, t+0.08, w-0.16, 0.3, size=12, bold=True, color=color)
    T(slide, desc, l+0.08, t+0.38, w-0.16, 0.4, size=10, color=GREY2)

# ══════════════════════════════════════════════════════════════════════════════
# Layout base de diapositiva de contenido
# ══════════════════════════════════════════════════════════════════════════════
# Columnas: izquierda bullets [0.4 → 7.0] | derecha visual [7.2 → 12.9]
BL = 0.4;  BW = 6.6   # bullets: left, width
VL = 7.2;  VW = 5.7   # visual:  left, width
CT = 1.15;  CH = 6.1   # content top, height

def slide_base(slide, title):
    R(slide, 0, 0, W, H, BG)
    R(slide, 0, 0, W, 0.08, ACCENT)
    R(slide, 0, 0.08, W, 0.97, CARD)
    T(slide, title, 0.5, 0.12, 11.5, 0.82, size=30, bold=True, color=ACCENT)
    # Panel bullets
    R(slide, BL, CT, BW, CH, CARD)
    R(slide, BL, CT, 0.06, CH, ACCENT)
    # Panel visual
    R(slide, VL, CT, VW, CH, CARD)

def bullets(slide, items, start_y=None):
    y = (start_y or CT) + 0.25
    for item in items:
        R(slide, BL+0.28, y+0.12, 0.10, 0.10, ACCENT)
        T(slide, item, BL+0.45, y, BW-0.55, 0.65, size=16, color=GREY)
        y += 0.75

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 1 — Portada
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
R(s, 0, 0, W, H, BG)
R(s, 0, 0, W, 0.12, ACCENT)
R(s, 0, H-0.12, W, 0.12, ACCENT)
R(s, 1.4, 1.6, 10.5, 4.3, CARD)
R(s, 1.4, 1.6, 0.1, 4.3, ACCENT)

# Decoración geométrica fondo
R(s, 9.5, 1.8, 3.5, 3.5, RGBColor(0x0A,0x22,0x35))
R(s, 10.0, 2.3, 2.5, 2.5, RGBColor(0x07,0x18,0x26))
R(s, 10.3, 0.9, 0.08, 5.8, RGBColor(0x00,0x80,0xA0))

T(s, "AndroidCommonDoc",
  2.0, 2.0, 9.0, 1.1, size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
T(s, "Kit centralizado para Android y Kotlin Multiplatform",
  2.0, 3.1, 9.0, 0.75, size=20, color=WHITE, align=PP_ALIGN.CENTER)
T(s, "Presentación de ingeniería  |  2026",
  2.0, 3.95, 9.0, 0.5, size=14, color=GREY2, align=PP_ALIGN.CENTER)
T(s, "Escalable  ·  Multiplataforma  ·  Nativo para IA",
  2.0, 4.55, 9.0, 0.5, size=13, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)
notes(s, "Introduce la iniciativa. AndroidCommonDoc es nuestro toolkit L0: fuente única de verdad para scripts, habilidades de IA, reglas de arquitectura y documentación. Instalas una vez y se usa en todos los proyectos.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 2 — Resumen ejecutivo  →  4 tarjetas KPI
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Resumen ejecutivo")
bullets(s, [
    "Fuente única de verdad para todos los estándares de desarrollo Android y KMP",
    "33 habilidades de IA compartidas entre Claude Code y GitHub Copilot",
    "13 reglas Detekt personalizadas aplicadas en tiempo de compilación",
    "17 herramientas MCP para validación programática y monitoreo",
    "Proyectos descendentes heredan mediante registro + manifiesto + sincronización",
])
# KPI cards 2x2
kpi_card(s, VL+0.2,  CT+0.25, 2.55, 1.25, "33", "Habilidades de IA", ACCENT)
kpi_card(s, VL+2.95, CT+0.25, 2.55, 1.25, "13", "Reglas Detekt",     GREEN)
kpi_card(s, VL+0.2,  CT+1.65, 2.55, 1.25, "17", "Herramientas MCP",  ORANGE)
kpi_card(s, VL+2.95, CT+1.65, 2.55, 1.25, "3",  "Capas L0/L1/L2",    PURPLE)
# Separador
R(s, VL+0.2, CT+3.1, VW-0.4, 0.02, ACCENT2)
T(s, "Mision", VL+0.2, CT+3.25, VW-0.4, 0.3, size=11, bold=True, color=ACCENT2)
T(s, "Instalar una vez · Usar en todos los proyectos · Sin duplicación",
  VL+0.2, CT+3.6, VW-0.4, 0.5, size=12, color=GREY, italic=True)
notes(s, "Resumen en 30 segundos: un repositorio con todos los patrones, scripts y habilidades. Proyectos heredan por sincronización. Las violaciones se detectan en compilación. Los agentes comparten las mismas habilidades.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 3 — El problema  →  Contraste antes/después
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "El problema")
bullets(s, [
    "Scripts duplicados en cada proyecto: misma corrección N veces",
    "Patrones inconsistentes generan fricción en revisiones y errores",
    "Los agentes de IA consumen hasta 50 000 tokens por ejecución",
    "Sin detección de cambios o deprecaciones en librerías externas",
    "Los patrones están documentados pero no se verifican automáticamente",
])
# Panel antes / después
T(s, "Sin toolkit", VL+0.2, CT+0.2, 2.5, 0.35, size=13, bold=True, color=RED)
for i, txt in enumerate([
    "Scripts distintos en cada proyecto",
    "50 000 tokens por ejecución de tests",
    "Patrones solo en documentos Word",
    "Errores descubiertos en producción",
]):
    R(s, VL+0.2, CT+0.6+i*0.82, 2.5, 0.65, RGBColor(0x2A,0x10,0x10), line_color=RED, line_w=0.8)
    T(s, txt, VL+0.3, CT+0.65+i*0.82, 2.3, 0.55, size=11, color=RGBColor(0xFF,0xAA,0xAA))

T(s, "Con toolkit", VL+2.95, CT+0.2, 2.55, 0.35, size=13, bold=True, color=GREEN)
for i, txt in enumerate([
    "Un script, todos los proyectos",
    "2 000 tokens (reducción x25)",
    "Reglas Detekt en compilación",
    "Violaciones detectadas en CI",
]):
    R(s, VL+2.95, CT+0.6+i*0.82, 2.55, 0.65, RGBColor(0x0A,0x2A,0x15), line_color=GREEN, line_w=0.8)
    T(s, txt, VL+3.05, CT+0.65+i*0.82, 2.35, 0.55, size=11, color=RGBColor(0x90,0xFF,0xB0))
notes(s, "Antes: cada proyecto tenía su propio script, los agentes leían 200+ líneas de Gradle, la deriva de versiones se descubría en producción.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 4 — Visión general  →  Hub diagram
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Visión general de la solución")
bullets(s, [
    "Toolkit L0 centralizado: instalar una vez, usar en todos los proyectos",
    "Scripts multiplataforma (PowerShell + Bash) con comportamiento idéntico",
    "Habilidades de IA en Markdown, invocables desde Claude Code o Copilot",
    "Reglas Detekt aplican la arquitectura en compilación (coste cero en runtime)",
    "Servidor MCP expone 17 herramientas para CI y automatización",
])
# Hub central
R(s, VL+1.7, CT+0.3, 2.3, 0.8, ACCENT, line_color=ACCENT2, line_w=1)
T(s, "AndroidCommonDoc", VL+1.7, CT+0.3, 2.3, 0.8, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
# Nodos satélite
nodes = [
    (VL+0.1, CT+1.5, "Scripts"),
    (VL+1.8, CT+1.5, "Habilidades IA"),
    (VL+3.5, CT+1.5, "Reglas Detekt"),
    (VL+0.6, CT+2.8, "Servidor MCP"),
    (VL+2.9, CT+2.8, "Docs & Vault"),
]
for nl, nt, nlabel in nodes:
    R(s, nl, nt, 1.85, 0.6, CARD2, line_color=ACCENT2, line_w=0.8)
    T(s, nlabel, nl, nt, 1.85, 0.6, size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
# Líneas conectoras (rectángulos finos)
for nl, nt, _ in nodes:
    cx = nl + 0.925; cy = nt + 0.3
    hc = VL + 1.7 + 1.15; vc = CT + 0.3 + 0.4
    lw = abs(cx - hc) if abs(cx - hc) > 0.05 else 0.02
    lh = abs(cy - vc) if abs(cy - vc) > 0.05 else 0.02
    R(s, min(cx,hc), min(cy,vc), lw, 0.02, RGBColor(0x33,0x66,0x80))
    R(s, min(cx,hc)+lw/2, min(cy,vc), 0.02, lh, RGBColor(0x33,0x66,0x80))
T(s, "L1 / L2 consumen via sync-l0", VL+0.2, CT+3.8, VW-0.4, 0.4, size=11, italic=True, color=GREY2, align=PP_ALIGN.CENTER)
notes(s, "La idea central: separar estructura de ejecución. Los scripts hacen el trabajo pesado. Las habilidades indican al agente qué hacer. Las reglas aplican decisiones. El servidor MCP hace el toolkit programable.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 5 — Arquitectura L0/L1/L2  →  Pirámide de capas
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Arquitectura del sistema: L0 / L1 / L2")
bullets(s, [
    "L0 (este repo): patrones universales, habilidades, agentes y reglas Detekt",
    "L1 (libs compartidas): autoridad de versiones, convenciones, sobreescrituras",
    "L2 (aplicaciones): habilidades específicas, agentes de dominio, overrides",
    "Distribución mediante registro SHA-256 + manifiesto + motor de sincronización",
    "Proyectos descendentes excluyen entradas individualmente: ausencia = exclusión",
])
# Pirámide de 3 capas (cajas apiladas, la mayor abajo)
tier_box(s, VL+0.3, CT+0.25, VW-0.6, 1.6, "L0  —  AndroidCommonDoc", "Patrones · Skills · Agentes · Reglas Detekt · MCP", ACCENT)
tier_box(s, VL+0.7, CT+2.05, VW-1.4, 1.4, "L1  —  Libs compartidas", "Versiones · Convenciones · Overrides L0", GREEN)
tier_box(s, VL+1.1, CT+3.65, VW-2.2, 1.3, "L2  —  App / proyecto", "Skills específicas · Agentes de dominio", ORANGE)
# Flechas
T(s, "↓  sync-l0", VL+2.2, CT+1.88, 2.0, 0.3, size=12, color=GREY2, align=PP_ALIGN.CENTER)
T(s, "↓  sync-l0", VL+2.2, CT+3.48, 2.0, 0.3, size=12, color=GREY2, align=PP_ALIGN.CENTER)
notes(s, "Modelo de capas. L0 define la línea base canónica. La sincronización es pull-based y verificada por hash. La configuración Detekt usa CompositeConfig: L1 sobreescribe solo lo que necesita.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 6 — Stack tecnológico  →  Grid de tecnologías
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Stack tecnológico")
bullets(s, [
    "Kotlin Multiplatform — Android · Desktop · iOS · macOS",
    "Detekt: reglas AST puras, sin resolución de tipos (evita bug 8882)",
    "Node.js TypeScript: servidor MCP, 17 herramientas, Vitest 653/653",
    "Claude Code + GitHub Copilot: 33 habilidades, 34 plantillas Copilot",
    "PowerShell + Bash: scripts con paridad verificada por tests",
])
# Grid 2x3 de tecnologías
tech = [
    ("Kotlin\nMultiplatform", ACCENT,  "Android · iOS · Desktop · macOS"),
    ("Detekt\nRules",         GREEN,   "13 reglas AST · L0/L1 config"),
    ("Node.js\nTypeScript",   ORANGE,  "MCP server · Vitest · 17 tools"),
    ("Claude Code\nCopilot",  PURPLE,  "33 skills · 34 templates"),
    ("Bash\nPowerShell",      YELLOW,  "Cross-platform · paridad verificada"),
    ("GitHub\nActions",       ACCENT2, "4 reusable workflows · CI/CD"),
]
cols = 2; rows = 3
cw = (VW - 0.4) / cols
ch = (CH - 0.5) / rows
for i, (name, color, sub) in enumerate(tech):
    col = i % cols; row = i // cols
    l = VL + 0.2 + col * cw
    t = CT + 0.25 + row * ch
    R(s, l, t, cw-0.15, ch-0.15, CARD2, line_color=color, line_w=1.2)
    R(s, l, t, cw-0.15, 0.05, color)
    T(s, name, l+0.1, t+0.1, cw-0.3, 0.55, size=13, bold=True, color=color)
    T(s, sub,  l+0.1, t+0.65, cw-0.3, 0.45, size=10, color=GREY2)
notes(s, "Las reglas Detekt solo usan AST para evitar un bug conocido de ClassCastException. El servidor MCP usa transporte stdio. Las habilidades son Markdown, sin dependencia de proveedor de IA.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 7 — Flujo de trabajo IA  →  Diagrama de secuencia
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Flujo de trabajo con habilidades de IA")
bullets(s, [
    "Desarrollador escribe /test core:domain en Claude Code o Copilot Chat",
    "El agente lee SKILL.md e invoca el script multiplataforma",
    "El script ejecuta Gradle, analiza XML y extrae errores estructurados",
    "El agente recibe ~2 000 tokens de señal (no ~50 000 de logs en bruto)",
    "La misma habilidad funciona en Windows (PowerShell) y macOS/Linux (Bash)",
])
# Diagrama de secuencia vertical
steps = [
    ("👤  Desarrollador",   "/test core:domain", ACCENT),
    ("🤖  Agente IA",       "Lee SKILL.md · invoca script", PURPLE),
    ("⚙️  Script",          "Gradle · XML → resumen", ORANGE),
    ("📊  Resultado",       "~2K tokens · tabla Markdown", GREEN),
]
bx = VL + 0.4
for i, (actor, action, color) in enumerate(steps):
    ty = CT + 0.3 + i * 1.35
    R(s, bx, ty, VW-0.8, 1.1, CARD2, line_color=color, line_w=1.5)
    R(s, bx, ty, 0.07, 1.1, color)
    T(s, actor,  bx+0.18, ty+0.08, VW-1.1, 0.38, size=13, bold=True, color=color)
    T(s, action, bx+0.18, ty+0.48, VW-1.1, 0.45, size=12, color=GREY)
    if i < len(steps)-1:
        T(s, "↓", bx + (VW-0.8)/2, ty+1.12, 0.4, 0.22, size=14, color=GREY2, align=PP_ALIGN.CENTER)
notes(s, "Sin el toolkit el agente lee 200+ líneas de Gradle. Con él recibe una tabla Markdown de resultados. La habilidad /coverage prohíbe explícitamente leer archivos XML. Reducción de tokens: ~25x.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 8 — Reglas Detekt  →  Tarjetas por categoría
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Reglas de arquitectura Detekt — 13 reglas")
bullets(s, [
    "Estado y exposición: SealedUiState · MutableStateFlowExposed · WhileSubscribedTimeout",
    "Límites del ViewModel: NoPlatformDeps · NoHardcodedStrings · NoHardcodedDispatchers",
    "Corrutinas: CancellationExceptionRethrow · NoRunCatching · NoSilentCatch · NoLaunchInInit",
    "Guardias de arquitectura: NoChannelForUiEvents · NoChannelForNavigation · NoMagicNumbers",
    "Configuración base L0 activa; L1 sobreescribe solo lo que necesita",
])
# 4 tarjetas de categoría
cats = [
    ("Estado y Exposición", "3 reglas", ACCENT,  ["SealedUiState", "MutableStateFlowExposed", "WhileSubscribedTimeout"]),
    ("Límites ViewModel",   "3 reglas", GREEN,   ["NoPlatformDeps", "NoHardcodedStrings", "NoHardcodedDispatchers"]),
    ("Corrutinas",          "4 reglas", ORANGE,  ["CancellationExceptionRethrow", "NoRunCatching", "NoSilentCatch", "NoLaunchInInit"]),
    ("Guardias Arquitectura","3 reglas", PURPLE, ["NoChannelForUiEvents", "NoChannelForNavigation", "NoMagicNumbers"]),
]
cw2 = (VW - 0.4) / 2
ch2 = (CH - 0.4) / 2
for i, (cat, count, color, rules) in enumerate(cats):
    col = i % 2; row = i // 2
    l2 = VL + 0.2 + col * cw2
    t2 = CT + 0.2 + row * ch2
    R(s, l2, t2, cw2-0.15, ch2-0.15, CARD2, line_color=color, line_w=1.2)
    R(s, l2, t2, cw2-0.15, 0.05, color)
    T(s, cat,   l2+0.12, t2+0.1,  cw2-0.3, 0.32, size=13, bold=True, color=color)
    T(s, count, l2+0.12, t2+0.42, cw2-0.3, 0.25, size=11, color=GREY2)
    for j, rule in enumerate(rules):
        R(s, l2+0.12, t2+0.7+j*0.37, 0.12, 0.12, color)
        T(s, rule, l2+0.3, t2+0.67+j*0.37, cw2-0.45, 0.32, size=10, color=GREY)
notes(s, "Todas las reglas son solo AST. Detectan antipatrones que superan la revisión de código. Cero falsos positivos: companion object queda correctamente excluido. Adopción en una línea con el plugin de convención.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 9 — Pipeline CI/CD  →  Etapas del pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Pipeline de CI/CD")
bullets(s, [
    "4 workflows reutilizables: commit-lint, recursos, seguridad KMP, guardias de arquitectura",
    "Hooks Claude Code: Detekt en tiempo real en cada escritura de archivo y commit",
    "quality-gate-orchestrator: informe unificado de 5 puertas de calidad",
    "Guardias Konsist: nomenclatura, límites de capas, aplicación de feature gates",
    "653/653 tests del servidor MCP en verde; cero fallos en la línea base",
])
# Pipeline horizontal: 4 etapas con flechas
stages = [
    ("Escritura\nde código", "Hook Detekt\nen tiempo real", ACCENT),
    ("Git commit", "Pre-commit\nvalidation hook", ORANGE),
    ("Pull Request", "Reusable\nworkflows CI", GREEN),
    ("Puerta\nde calidad", "5 gates\nPASS / FAIL", PURPLE),
]
sw = (VW - 0.3) / len(stages) - 0.05
for i, (stage, desc, color) in enumerate(stages):
    l3 = VL + 0.2 + i * (sw + 0.3)
    R(s, l3, CT+0.3, sw, 1.5, CARD2, line_color=color, line_w=1.5)
    R(s, l3, CT+0.3, sw, 0.06, color)
    T(s, stage, l3, CT+0.4, sw, 0.55, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, desc,  l3, CT+0.95, sw, 0.65, size=11, color=GREY, align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        T(s, "→", l3+sw+0.01, CT+0.8, 0.28, 0.4, size=18, bold=True, color=GREY2, align=PP_ALIGN.CENTER)

# Detalle: lista de reusable workflows
R(s, VL+0.2, CT+2.2, VW-0.4, 0.04, ACCENT2)
T(s, "Reusable workflows incluidos:", VL+0.2, CT+2.35, VW-0.4, 0.3, size=11, bold=True, color=ACCENT2)
wf_items = ["reusable-commit-lint.yml", "reusable-lint-resources.yml", "reusable-kmp-safety-check.yml", "reusable-architecture-guards.yml"]
for j, wf in enumerate(wf_items):
    col3 = j % 2; row3 = j // 2
    T(s, "▸ " + wf, VL+0.2+col3*(VW/2-0.2), CT+2.75+row3*0.45, VW/2-0.2, 0.38, size=11, color=GREEN)
notes(s, "Los workflows reutilizables permiten que cada app obtenga CI importando un archivo. Los hooks son shift-left: violaciones antes del commit. El orchestrator ejecuta las 5 puertas y produce un único PASS/FAIL.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 10 — Estrategia de pruebas  →  Pirámide de testing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Estrategia de pruebas")
bullets(s, [
    "Unitarias: Vitest 653/653 para MCP · Gradle 76/76 para reglas Detekt",
    "Estructurales: Konsist verifica nomenclatura, límites de capas y feature gates",
    "Cobertura: Kover (KMP) o JaCoCo con ejecución paralela y análisis de lagunas",
    "Integración: estructura de docs, validación CLAUDE.md e integridad del vault",
    "Habilidad /pre-pr: lint + pruebas + formato de commit antes de cada merge",
])
# Pirámide invertida (más estrechas arriba = menos tests, base = más tests)
pyramid = [
    ("Unitarias",     "653/653 MCP · 76/76 Detekt",   ACCENT,  4.5),
    ("Integración",   "doc-structure · CLAUDE.md · vault", GREEN, 3.5),
    ("Estructurales", "Konsist guards · 17 tests",     ORANGE,  2.5),
    ("Manual / E2E",  "/pre-pr skill antes de merge",  PURPLE,  1.5),
]
cx4 = VL + VW/2
for i, (name, desc, color, width) in enumerate(pyramid):
    ty4 = CT + 0.25 + i * 1.2
    lx4 = cx4 - width/2
    R(s, lx4, ty4, width, 0.95, CARD2, line_color=color, line_w=1.5)
    R(s, lx4, ty4, 0.06, 0.95, color)
    T(s, name, lx4+0.15, ty4+0.08, width-0.2, 0.35, size=13, bold=True, color=color)
    T(s, desc, lx4+0.15, ty4+0.45, width-0.2, 0.38, size=11, color=GREY)
notes(s, "Tests unitarios para toda la lógica de negocio. Los estructurales detectan deriva arquitectónica antes de que sea deuda técnica. Las herramientas de cobertura identifican líneas no probadas y generan esqueletos. /pre-pr orquesta todo antes del PR.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 11 — Servidor MCP  →  Catálogo de herramientas
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Servidor MCP e inteligencia documental")
bullets(s, [
    "17 herramientas sobre transporte stdio; integra con Claude Desktop y CI",
    "validate-all: ejecuta todos los scripts de validación → JSON estructurado",
    "monitor-sources: monitoreo por niveles (GitHub API · Maven Central · docs)",
    "Detección de versión: cuerpo del release analizado solo cuando hay cambio",
    "Cambios en páginas de docs: hash SHA-256, sin escaneo de palabras clave",
])
# Grupos de herramientas
groups = [
    ("Validación", ACCENT,  ["validate-all","verify-kmp","check-version-sync","validate-skills","validate-doc-structure","validate-vault"]),
    ("Monitoreo",  GREEN,   ["monitor-sources","check-doc-freshness"]),
    ("Generación", ORANGE,  ["generate-detekt-rules","ingest-content"]),
    ("Vault",      PURPLE,  ["sync-vault","vault-status","validate-vault"]),
    ("Utilidad",   YELLOW,  ["find-pattern","setup-check","rate-limit-status","script-parity"]),
]
gw = (VW - 0.3) / 2 - 0.1
gy = CT + 0.2
for i, (gname, gcolor, tools) in enumerate(groups):
    col5 = i % 2; row5 = i // 2
    gl = VL + 0.2 + col5 * (gw + 0.2)
    gt = gy + row5 * 2.05
    R(s, gl, gt, gw, 1.85, CARD2, line_color=gcolor, line_w=1)
    R(s, gl, gt, gw, 0.05, gcolor)
    T(s, gname, gl+0.1, gt+0.1, gw-0.2, 0.3, size=12, bold=True, color=gcolor)
    for j, tool in enumerate(tools[:4]):
        T(s, "· "+tool, gl+0.1, gt+0.45+j*0.33, gw-0.2, 0.28, size=10, color=GREY)
notes(s, "Transporte stdio: sin servidor de red que mantener. Rate limited a 30 llamadas/min. monitor-sources es basado en señales: solo analiza notas del release cuando la versión cambia. Los cambios en docs se detectan por hash.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 12 — Impacto en el negocio  →  Gráfico de barras
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Impacto en el negocio")
bullets(s, [
    "Reducción de ~25× en tokens por ejecución de pruebas asistida por IA",
    "Violaciones de arquitectura en compilación: cero ciclos extra de revisión",
    "Nuevo proyecto configurado en minutos, no en semanas",
    "Deprecaciones de librerías detectadas automáticamente antes del lanzamiento",
    "Cada nueva app hereda el pipeline CI completo con una sola importación",
])
# Gráfico de barras horizontal comparativo
T(s, "Coste de tokens por ejecución de tests", VL+0.2, CT+0.25, VW-0.4, 0.35, size=13, bold=True, color=ACCENT2)
R(s, VL+0.2, CT+0.62, VW-0.4, 0.02, GREY2)

metrics = [
    ("Sin toolkit",  "~50 000 tokens", RED,    1.0),
    ("Con toolkit",  "~2 000 tokens",  GREEN,  0.04),
]
bar_w = VW - 0.8
for i, (label, val, color, pct) in enumerate(metrics):
    ty5 = CT + 0.8 + i * 0.9
    T(s, label, VL+0.2, ty5, 2.2, 0.28, size=12, color=GREY)
    R(s, VL+0.2, ty5+0.32, bar_w, 0.3, CARD2)
    R(s, VL+0.2, ty5+0.32, max(bar_w*pct, 0.1), 0.3, color)
    T(s, val, VL+0.2+max(bar_w*pct,0.1)+0.1, ty5+0.3, 1.8, 0.32, size=12, bold=True, color=color)

R(s, VL+0.2, CT+2.85, VW-0.4, 0.02, GREY2)
T(s, "Otros indicadores", VL+0.2, CT+3.0, VW-0.4, 0.3, size=12, bold=True, color=ACCENT2)
kpis = [
    ("Onboarding", "< 15 min por proyecto", GREEN),
    ("Violaciones en CI", "~30 min ahorrados por violación", ORANGE),
    ("Duplicación", "Scripts compartidos entre todos los proyectos", ACCENT),
]
for i, (k, v, c) in enumerate(kpis):
    R(s, VL+0.2, CT+3.4+i*0.7, 0.12, 0.12, c)
    T(s, f"{k}: {v}", VL+0.42, CT+3.35+i*0.7, VW-0.6, 0.35, size=11, color=GREY)
notes(s, "ROI: 20 ejecuciones al día suponen ~1M tokens/mes ahorrados por desarrollador. Cada violación en CI ahorra ~30 min de ronda de PR. Nuevos proyectos de días a minutos. Sin sorpresas en el lanzamiento.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 13 — Casos de uso reales  →  Comandos en acción
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Casos de uso reales")
bullets(s, [
    "/pre-pr: lint · pruebas · formato de commit · tabla resumen antes del PR",
    "/auto-cover: genera pruebas para líneas sin cobertura automáticamente",
    "/sync-l0: propagó 4 nuevas habilidades a DawSync en menos de un minuto",
    "Hook Detekt bloqueó un commit con CancellationException silenciada (Copilot)",
    "quality-gate-orchestrator encontró 3 contadores desactualizados en README",
])
# Tarjetas de comandos
cmds = [
    ("/pre-pr",     "Orquesta lint + tests + commit-lint\nantes de cada merge",  ACCENT),
    ("/auto-cover", "Genera tests para líneas\nsin cobertura identificadas",     GREEN),
    ("/sync-l0",    "Propaga habilidades L0\na proyectos descendentes",          ORANGE),
    ("/sbom-scan",  "Genera SBOM + escanea\nCVEs con Trivy",                     PURPLE),
]
cw6 = (VW - 0.3) / 2 - 0.08
ch6 = (CH - 0.5) / 2 - 0.1
for i, (cmd, desc, color) in enumerate(cmds):
    col6 = i % 2; row6 = i // 2
    l6 = VL + 0.2 + col6 * (cw6 + 0.2)
    t6 = CT + 0.25 + row6 * (ch6 + 0.2)
    R(s, l6, t6, cw6, ch6, CARD2, line_color=color, line_w=1.5)
    R(s, l6, t6, cw6, 0.06, color)
    T(s, cmd,  l6+0.12, t6+0.12, cw6-0.2, 0.4,  size=16, bold=True, color=color)
    T(s, desc, l6+0.12, t6+0.55, cw6-0.2, 0.85, size=12, color=GREY)
notes(s, "/pre-pr impidió PRs incompletos. /auto-cover identificó cobertura perdida en UseCases. sync-l0 propagó habilidades en menos de un minuto. El hook detectó CancellationException silenciada por autocompletado de Copilot.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 14 — Seguridad y rendimiento  →  Checklist visual
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Seguridad y rendimiento")
bullets(s, [
    "Generación de SBOM (CycloneDX) y análisis de CVE (Trivy) via /sbom-scan",
    "Konsist: FeatureAccess.GRANTED solo lo devuelven implementaciones autorizadas",
    "Sin secretos en el código: Detekt NoHardcodedStrings + Konsist DS-PROD1",
    "Ejecución paralela de pruebas: una invocación Gradle, 2-3× más rápido",
    "Rate limiting servidor MCP: 30 llamadas/min protege APIs externas",
])
# Panel de seguridad dividido en dos zonas
T(s, "Seguridad", VL+0.3, CT+0.2, 2.5, 0.32, size=13, bold=True, color=RED)
sec_items = [
    ("SBOM generado en cada release",     GREEN),
    ("Sin secretos en código",            GREEN),
    ("Feature gate es obligatorio",       GREEN),
    ("CVE scan en CI con Trivy",          GREEN),
    ("NoGlobalScope: sin fugas de scope", GREEN),
]
for i, (item, color) in enumerate(sec_items):
    R(s, VL+0.3, CT+0.65+i*0.62, 0.28, 0.28, color)
    T(s, "✓", VL+0.31, CT+0.62+i*0.62, 0.26, 0.3, size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, item, VL+0.7, CT+0.62+i*0.62, 2.5, 0.32, size=11, color=GREY)

R(s, VL+3.2, CT+0.2, 0.03, CH-0.4, GREY2)

T(s, "Rendimiento", VL+3.4, CT+0.2, 2.2, 0.32, size=13, bold=True, color=ORANGE)
perf_items = [
    ("Tokens: 50K → 2K (×25)",    ORANGE),
    ("Tests paralelos: 2-3× más rápido", ORANGE),
    ("Rate limiting: 30 calls/min", ORANGE),
    ("CI reutilizable: 1 archivo",  ORANGE),
]
for i, (item, color) in enumerate(perf_items):
    R(s, VL+3.4, CT+0.65+i*0.75, 0.28, 0.28, color)
    T(s, "↑", VL+3.41, CT+0.62+i*0.75, 0.26, 0.3, size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, item, VL+3.8, CT+0.62+i*0.75, 2.0, 0.32, size=11, color=GREY)
notes(s, "SBOM mantiene inventario de software para cada release. El guardia de feature gate es nivel crítico: previene bypass de facturación. NoGlobalScope previene fugas de corrutinas. Paralelo reduce minutos de CI en GitHub Actions.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 15 — Adopción  →  5 pasos numerados
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Adopción: primeros pasos")
bullets(s, [
    "1. Clonar el repositorio y configurar ANDROID_COMMON_DOC",
    "2. Ejecutar setup-toolkit.sh --project-root /tu/proyecto",
    "3. Añadir l0-manifest.json y declarar entradas L0",
    "4. Ejecutar /sync-l0 para materializar las habilidades",
    "5. Añadir id(\"androidcommondoc.toolkit\") al plugin de Gradle",
])
# 5 pasos con número grande
step_h = (CH - 0.4) / 5
for i, (num, label, sublabel) in enumerate([
    ("1", "git clone + export ANDROID_COMMON_DOC=...",   "Configura la variable de entorno"),
    ("2", "bash setup-toolkit.sh --project-root /ruta",  "Instala hooks, plugin y Copilot instructions"),
    ("3", "Crea l0-manifest.json",                       "Declara qué habilidades L0 sincronizar"),
    ("4", "/sync-l0",                                    "Materializa habilidades con tracking de versión"),
    ("5", "id(\"androidcommondoc.toolkit\")",            "Activa las 13 reglas Detekt en compilación"),
]):
    ty7 = CT + 0.2 + i * step_h
    R(s, VL+0.2, ty7, VW-0.4, step_h-0.1, CARD2, line_color=ACCENT, line_w=0.8)
    R(s, VL+0.2, ty7, 0.55, step_h-0.1, ACCENT)
    T(s, num,      VL+0.21, ty7+0.05, 0.53, step_h-0.2, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, label,    VL+0.85, ty7+0.05, VW-1.2, 0.35, size=12, bold=True, color=WHITE)
    T(s, sublabel, VL+0.85, ty7+0.43, VW-1.2, 0.3,  size=11, color=GREY2)
T(s, "⏱  < 15 minutos para un proyecto nuevo", VL+0.2, CT+CH-0.25, VW-0.4, 0.28,
  size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
notes(s, "Proceso de incorporación: 5 pasos, menos de 15 minutos para un proyecto nuevo. El setup-toolkit.sh instala hooks, instrucciones de Copilot y el plugin de convención automáticamente.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 16 — Hoja de ruta  →  Timeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Hoja de ruta")
bullets(s, [
    "M008: emitBaselineConfig → MCP generate-detekt-rules (regeneración automática)",
    "M008: workflows CI reutilizables adoptados en DawSync (pre-pr, git-flow, commit-lint)",
    "Monitoreo v2: diferencias de changelog por versión + referencias cruzadas de API",
    "Plantillas Konsist ampliadas con más contratos arquitectónicos de DawSync",
    "Explorar distribución como plugin Gradle publicado en Maven Central",
])
# Timeline horizontal
phases = [
    ("M007\n✓ Hecho",     "13 reglas Detekt\nCI workflows\nGit Flow skills",    GREEN,  "Completado"),
    ("M008\n→ Próximo",   "emitBaselineConfig\nDawSync CI adoption\nMonitoreo v2", ORANGE, "En curso"),
    ("M009\n◦ Futuro",    "Konsist plantillas\nMaven Central\nPlugin público",   GREY2,  "Planificado"),
]
pw = (VW - 0.4) / len(phases)
# Línea de tiempo
R(s, VL+0.2, CT+1.0, VW-0.4, 0.04, GREY2)
for i, (phase, desc, color, status) in enumerate(phases):
    lp = VL + 0.2 + i * pw
    # Punto en línea
    R(s, lp + pw/2 - 0.12, CT+0.85, 0.24, 0.24, color)
    # Tarjeta
    R(s, lp+0.1, CT+1.25, pw-0.2, 4.4, CARD2, line_color=color, line_w=1.2)
    R(s, lp+0.1, CT+1.25, pw-0.2, 0.05, color)
    T(s, phase,  lp+0.2, CT+1.35, pw-0.4, 0.7, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, status, lp+0.2, CT+2.08, pw-0.4, 0.3, size=10, italic=True, color=GREY2, align=PP_ALIGN.CENTER)
    R(s, lp+0.1, CT+2.45, pw-0.2, 0.03, color)
    T(s, desc,   lp+0.2, CT+2.6,  pw-0.4, 1.6, size=11, color=GREY, align=PP_ALIGN.CENTER)
notes(s, "Prioridades a corto plazo: emitBaselineConfig conectado al MCP. La adopción de workflows en DawSync es el primer adoptante L2 de los workflows reutilizables. Monitoreo v2: señal de calidad frente al ruido de palabras clave.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 17 — Próximos pasos  →  Llamada a la acción
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_base(s, "Próximos pasos")
bullets(s, [
    "Prueba /pre-pr en tu proyecto actual hoy mismo",
    "Añade l0-manifest.json y ejecuta /sync-l0 (33 habilidades disponibles)",
    "Revisa el catálogo de reglas Detekt: identifica las relevantes para tu código",
    "Conecta Claude Desktop al servidor MCP para validación programática",
    "Consulta docs/guides/detekt-config.md o abre un issue en el repositorio",
])
# Panel de llamada a la acción
T(s, "¿Por dónde empezar?", VL+0.3, CT+0.2, VW-0.5, 0.35, size=14, bold=True, color=ACCENT2)
R(s, VL+0.2, CT+0.58, VW-0.4, 0.03, ACCENT)
actions = [
    ("Hoy",      "/pre-pr",     "Sin configuración previa",             ACCENT),
    ("Semana 1", "/sync-l0",    "l0-manifest.json + 33 habilidades",    GREEN),
    ("Semana 2", "Detekt rules","Plugin de convención + 13 reglas",      ORANGE),
    ("Semana 3", "MCP server",  "npm start + configurar Claude Desktop", PURPLE),
]
for i, (when, cmd, desc, color) in enumerate(actions):
    ty8 = CT + 0.8 + i * 1.2
    R(s, VL+0.2, ty8, 1.1, 1.0, CARD2, line_color=color, line_w=1.2)
    T(s, when, VL+0.2, ty8+0.08, 1.1, 0.38, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    R(s, VL+1.45, ty8+0.05, VW-1.65, 0.9, CARD2, line_color=color, line_w=0.8)
    T(s, cmd,  VL+1.6, ty8+0.08, VW-1.85, 0.35, size=13, bold=True, color=color)
    T(s, desc, VL+1.6, ty8+0.47, VW-1.85, 0.35, size=11, color=GREY)
notes(s, "/pre-pr funciona inmediatamente sin configuración. sync-l0 es la puerta de entrada al toolkit completo. Reglas Detekt requieren plugin de convención (10 min). MCP server: npm start en mcp-server/.")

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 18 — Gracias
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
R(s, 0, 0, W, H, BG)
R(s, 0, 0, W, 0.12, ACCENT)
R(s, 0, H-0.12, W, 0.12, ACCENT)
# Decoración
for i in range(6):
    R(s, 9.0+i*0.65, 0.5+i*0.5, 4.0, 0.04, RGBColor(0x00,0x50,0x70))
R(s, 2.0, 1.6, 9.3, 4.5, CARD)
R(s, 2.0, 1.6, 0.1, 4.5, ACCENT)
T(s, "Gracias",
  2.5, 2.0, 8.5, 1.1, size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
T(s, "github.com/<org>/AndroidCommonDoc",
  2.5, 3.2, 8.5, 0.6, size=18, color=ACCENT2, align=PP_ALIGN.CENTER)
T(s, "docs/  ·  skills/  ·  mcp-server/  ·  detekt-rules/",
  2.5, 3.95, 8.5, 0.5, size=14, color=GREY, align=PP_ALIGN.CENTER)
R(s, 3.0, 4.62, 7.3, 0.03, GREY2)
T(s, "Preguntas bienvenidas",
  2.5, 4.75, 8.5, 0.5, size=14, italic=True, color=GREY2, align=PP_ALIGN.CENTER)
# Mini stats en la parte inferior
for i, (val, lbl, c) in enumerate([("33","Habilidades",ACCENT),("13","Reglas Detekt",GREEN),("17","Herramientas MCP",ORANGE),("653/653","Tests en verde",PURPLE)]):
    R(s, 2.4+i*2.15, 5.5, 1.9, 0.9, CARD2, line_color=c, line_w=1)
    T(s, val, 2.4+i*2.15, 5.52, 1.9, 0.42, size=18, bold=True, color=c, align=PP_ALIGN.CENTER)
    T(s, lbl, 2.4+i*2.15, 5.94, 1.9, 0.3, size=10, color=GREY2, align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "Dejar visible durante las preguntas. URL pública: compartir en el chat.\n"
    "Mensaje clave: el toolkit está listo para producción, probado, y adoptarlo lleva menos de 15 minutos."
)

# ── Guardar ───────────────────────────────────────────────────────────────────
out = "AndroidCommonDoc-Presentation.pptx"
prs.save(out)
print(f"Guardado: {out}  ({len(prs.slides)} diapositivas)")
